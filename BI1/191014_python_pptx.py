import time
from pptx import Presentation
from pptx.util import Inches

imgdir="/media/desktop-bi-16/D2C67EE7C67ECAED/BI/Analysis/191014/qc_CMC018_019/"
img1=imgdir+'fastqc_per_base_sequence_quality_plot.png'
img2=imgdir+'fastqc_per_sequence_quality_scores_plot.png'
img3=imgdir+''
img4=imgdir+'fastqc_per_sequence_gc_content_plot.png'
img5=imgdir+'fastqc_per_base_n_content_plot.png'
img6=imgdir+'fastqc_sequence_duplication_levels_plot.png'
img7=imgdir+'fastqc_overrepresented_sequencesi_plot.png'
img8=imgdir+'fastqc_adapter_content_plot.png'

img0=imgdir+'fastqc_sequence_counts_plot_percentage.png'
#cutadapt_plot.png
#salmon_plot.png
#fastqc_sequence_counts_plot.png
#fastqc_sequence_counts_plot_percentage.png



prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_slide_layout)

#12.34=5inch 8.23=4inch

left = top = Inches(0.3)
pic1_1 = slide1.shapes.add_picture(img1, left, top, height=Inches(4))
pic1_2 = slide1.shapes.add_picture(img2, Inches(6), top, height=Inches(4))



slide2 = prs.slides.add_slide(blank_slide_layout)
pic2_1 = slide2.shapes.add_picture(img4, left, top, height=Inches(4))
pic2_2 = slide2.shapes.add_picture(img5, Inches(6), top, height=Inches(4))

slide3 = prs.slides.add_slide(blank_slide_layout)
pic3_1 = slide3.shapes.add_picture(img6, left, top, height=Inches(4))
pic3_2 = slide3.shapes.add_picture(img7, Inches(6), top, height=Inches(4))

slide4 = prs.slides.add_slide(blank_slide_layout)
pic4_1 = slide4.shapes.add_picture(img8, left, top, height=Inches(4))
#pic4_2 = slide4.shapes.add_picture(img9, Inches(6), top, height=Inches(4), weight=Inches(5))

prs.save('%s/QC%s.pptx'%(imgdir,time.strftime("%y%m%d")))

"""
left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
"""

